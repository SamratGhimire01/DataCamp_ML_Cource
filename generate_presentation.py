"""
Vertical and Horizontal Alignments in Curriculum
Ultra-modern, minimalist academic presentation — python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ─── Colour palette ────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0F, 0x20, 0x27)   # Deep Navy Blue
TEAL    = RGBColor(0x00, 0x80, 0x80)   # Elegant Teal
SLATE   = RGBColor(0x49, 0x50, 0x57)   # Slate Gray
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF8, 0xF9, 0xFA)   # Ultra-light gray
DARK    = RGBColor(0x21, 0x25, 0x29)   # Dark Charcoal
MID     = RGBColor(0x6C, 0x75, 0x7D)   # Mid gray for captions
RULE    = RGBColor(0xDE, 0xE2, 0xE6)   # Light rule line
TEAL_L  = RGBColor(0xE0, 0xF2, 0xF2)   # Teal tint background

# ─── Slide size (16:9 widescreen) ──────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

# ─── Helper: blank slide ───────────────────────────────────────────────────────
def blank_slide(prs):
    layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(layout)

# ─── Helper: set slide background ─────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

# ─── Helper: add text box ─────────────────────────────────────────────────────
def txb(slide, text, x, y, w, h,
        size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
        italic=False, font="Calibri", wrap=True, line_space=None, space_before=None):
    color = color or DARK
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font
    if line_space:
        from pptx.util import Pt as Ptx
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_space * 100)))
    return tf

# ─── Helper: add rect shape with optional fill / border ───────────────────────
def rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0.75)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    sf = shape.fill
    if fill:
        sf.solid(); sf.fore_color.rgb = fill
    else:
        sf.background()
    ln = shape.line
    if line_color:
        ln.color.rgb = line_color
        ln.width     = line_width
    else:
        ln.fill.background()
    return shape

# ─── Helper: multi-line text in a textbox ─────────────────────────────────────
def multi_txb(slide, lines, x, y, w, h,
              size=14, color=None, font="Calibri", bold_first=False,
              align=PP_ALIGN.LEFT, space_after=Pt(4)):
    color = color or DARK
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        pPr = p._p.get_or_add_pPr()
        spcAft = etree.SubElement(pPr, qn('a:spcAft'))
        spcPts = etree.SubElement(spcAft, qn('a:spcPts'))
        spcPts.set('val', str(int(space_after.pt * 100)))
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.color.rgb = color
        run.font.name  = font
        run.font.bold  = (bold_first and i == 0)
    return tf

# ─── Helper: add connector line (just a thin rect) ────────────────────────────
def hline(slide, x, y, w, color=RULE, height=Pt(1)):
    r = rect(slide, x, y, w, height, fill=color)
    r.line.fill.background()
    return r

def vline(slide, x, y, h, color=RULE, width=Pt(1)):
    r = rect(slide, x, y, width, h, fill=color)
    r.line.fill.background()
    return r

# ─── Helper: arrow connector (horizontal, pointing right) ─────────────────────
def arrow_right(slide, x, y, w, color=TEAL):
    """Draw a thin horizontal arrow pointing right."""
    # shaft
    shaft_h = Inches(0.04)
    shaft_y = y - shaft_h / 2
    rect(slide, x, shaft_y, w - Inches(0.12), shaft_h, fill=color)
    # arrowhead (small triangle via narrow rect)
    head_w = Inches(0.12)
    head_h = Inches(0.18)
    head_x = x + w - head_w
    head_y = y - head_h / 2
    # use a right-triangle approximation: stacked thin rects
    steps = 5
    for i in range(steps):
        rh = head_h * (steps - i) / steps
        rx = head_x + head_w * i / steps
        rw = head_w / steps
        ry = y - rh / 2
        rect(slide, rx, ry, rw, rh, fill=color)

def arrow_down(slide, x, y, h, color=TEAL):
    """Draw a thin vertical arrow pointing down."""
    shaft_w = Inches(0.04)
    shaft_x = x - shaft_w / 2
    rect(slide, shaft_x, y, shaft_w, h - Inches(0.12), fill=color)
    head_h = Inches(0.12)
    head_w = Inches(0.18)
    head_y = y + h - head_h
    steps = 5
    for i in range(steps):
        rw = head_w * (steps - i) / steps
        ry = head_y + head_h * i / steps
        rh = head_h / steps
        rx = x - rw / 2
        rect(slide, rx, ry, rw, rh, fill=color)

# ──────────────────────────────────────────────────────────────────────────────
#  SLIDE BUILDERS
# ──────────────────────────────────────────────────────────────────────────────

def slide1_title(prs):
    """Title slide — deep navy background, white text, elegant minimal."""
    s = blank_slide(prs)
    set_bg(s, NAVY)

    # Vertical centre block
    # Teal accent bar on left
    rect(s, Inches(0.9), Inches(2.0), Inches(0.08), Inches(3.5), fill=TEAL)

    # Main title
    txb(s, "Vertical and Horizontal\nAlignments in Curriculum",
        Inches(1.2), Inches(2.0), Inches(11.0), Inches(2.2),
        size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri")

    # Rule
    rect(s, Inches(1.2), Inches(4.35), Inches(4.5), Inches(0.04), fill=TEAL)

    # Author block
    txb(s, "Dhruba Prasad Niure, PhD",
        Inches(1.2), Inches(4.55), Inches(8), Inches(0.55),
        size=20, bold=True, color=WHITE, font="Calibri")
    txb(s, "Associate Professor, Tribhuvan University",
        Inches(1.2), Inches(5.1), Inches(8), Inches(0.45),
        size=16, color=RGBColor(0xAD, 0xB5, 0xBD), font="Calibri")


def slide2_curriculum(prs):
    """4-column curriculum framework."""
    s = blank_slide(prs)
    set_bg(s, WHITE)

    # Header
    txb(s, "Curriculum: A Plan for Educational Program",
        Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.75),
        size=32, bold=True, color=NAVY, align=PP_ALIGN.LEFT, font="Calibri")
    hline(s, Inches(0.6), Inches(1.22), Inches(12.1), color=TEAL, height=Inches(0.045))

    cols = [
        ("Why to Teach",  "Objectives",
         "Formulating competencies and learning outcomes to be achieved."),
        ("What to Teach", "Content",
         "Identifying the essential knowledge, values, and skills students should learn."),
        ("How to Teach",  "Methods",
         "Selecting effective and relevant teaching methods and activities."),
        ("How to Assess", "Evaluation",
         "Determining how to measure students' LOTS and HOTS."),
    ]

    col_w = Inches(2.9)
    gap   = Inches(0.3)
    start_x = Inches(0.7)
    card_y = Inches(1.5)
    card_h = Inches(5.4)

    for i, (question, label, desc) in enumerate(cols):
        cx = start_x + i * (col_w + gap)

        # card bg
        rect(s, cx, card_y, col_w, card_h, fill=LIGHT,
             line_color=RULE, line_width=Pt(0.5))

        # Teal top strip
        rect(s, cx, card_y, col_w, Inches(0.08), fill=TEAL)

        # Number circle (simple)
        num_x = cx + Inches(0.18)
        num_y = card_y + Inches(0.22)
        rect(s, num_x, num_y, Inches(0.44), Inches(0.44), fill=TEAL)
        txb(s, str(i + 1), num_x, num_y, Inches(0.44), Inches(0.44),
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

        # Question heading
        txb(s, question, cx + Inches(0.18), card_y + Inches(0.78), col_w - Inches(0.36), Inches(0.65),
            size=15, bold=True, color=NAVY, font="Calibri")

        # Label (teal)
        txb(s, label.upper(), cx + Inches(0.18), card_y + Inches(1.5), col_w - Inches(0.36), Inches(0.45),
            size=11, bold=True, color=TEAL, font="Calibri",
            align=PP_ALIGN.LEFT)

        # Divider
        hline(s, cx + Inches(0.18), card_y + Inches(2.0), col_w - Inches(0.36),
              color=RULE, height=Inches(0.025))

        # Description
        txb(s, desc, cx + Inches(0.18), card_y + Inches(2.15), col_w - Inches(0.36), Inches(2.8),
            size=13, color=SLATE, font="Calibri")


def slide3_alignment(prs):
    """2x2 grid showing vertical & horizontal alignment Economics/Sociology Grades 11-12."""
    s = blank_slide(prs)
    set_bg(s, WHITE)

    txb(s, "Vertical and Horizontal Alignments",
        Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.7),
        size=32, bold=True, color=NAVY, font="Calibri")
    txb(s, "Alignment between Economics and Sociology at Grades 11–12",
        Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.45),
        size=15, color=SLATE, font="Calibri")

    # Grid layout
    cell_w = Inches(3.4)
    cell_h = Inches(1.55)
    gap    = Inches(0.5)
    arrow  = Inches(0.85)

    # Row positions (Grade row label + 2 subject cells)
    rows = [("Grade 12", Inches(2.0)), ("Grade 11", Inches(4.75))]
    cols = [("Economics", Inches(2.4)), ("Sociology", Inches(6.75))]

    # Grade labels (left)
    for grade, gy in rows:
        rect(s, Inches(0.6), gy, Inches(1.5), cell_h,
             fill=NAVY, line_color=None)
        txb(s, grade, Inches(0.6), gy, Inches(1.5), cell_h,
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    # Subject header (top)
    for subj, cx in cols:
        rect(s, cx, Inches(1.55), cell_w, Inches(0.45),
             fill=TEAL, line_color=None)
        txb(s, subj, cx, Inches(1.55), cell_w, Inches(0.45),
            size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    # Cells
    for grade, gy in rows:
        for subj, cx in cols:
            rect(s, cx, gy, cell_w, cell_h,
                 fill=LIGHT, line_color=RULE, line_width=Pt(1))
            label = f"{subj}\n{grade}"
            txb(s, label, cx, gy, cell_w, cell_h,
                size=14, color=NAVY, align=PP_ALIGN.CENTER,
                bold=True, font="Calibri")

    # Horizontal arrows (↔ same grade, between subjects)
    for grade, gy in rows:
        mid_y = gy + cell_h / 2
        # right arrow: Econ → Soc
        arrow_right(s, cols[0][1] + cell_w + Inches(0.05), mid_y,
                    cols[1][1] - cols[0][1] - cell_w - Inches(0.05), color=TEAL)

    # Vertical arrows (↕ same subject, between grades)
    for subj, cx in cols:
        mid_x = cx + cell_w / 2
        # down arrow: Grade12 → Grade11
        arrow_down(s, mid_x, rows[0][1] + cell_h + Inches(0.05),
                   rows[1][1] - rows[0][1] - cell_h - Inches(0.05), color=NAVY)

    # Legend
    ly = Inches(6.6)
    rect(s, Inches(2.4), ly, Inches(0.25), Inches(0.25), fill=TEAL)
    txb(s, "Horizontal Organization: Scope & Integration",
        Inches(2.75), ly, Inches(4.0), Inches(0.35), size=12, color=SLATE, font="Calibri")
    rect(s, Inches(7.0), ly, Inches(0.25), Inches(0.25), fill=NAVY)
    txb(s, "Vertical Organization: Continuity & Sequence",
        Inches(7.35), ly, Inches(4.0), Inches(0.35), size=12, color=SLATE, font="Calibri")


def slide4_questions(prs):
    """4 key questions — clean, spacious."""
    s = blank_slide(prs)
    set_bg(s, WHITE)

    txb(s, "Questions to Be Answered",
        Inches(0.7), Inches(0.4), Inches(11.8), Inches(0.7),
        size=32, bold=True, color=NAVY, font="Calibri")

    questions = [
        ("01", "Continuity",   "What does continuity refer to in curriculum organization?"),
        ("02", "Sequence",     "What do you mean by sequence in curriculum design?"),
        ("03", "Scope",        "How do you define scope within the curriculum framework?"),
        ("04", "Integration",  "What does integration mean in the context of curriculum?"),
    ]

    q_y = Inches(1.45)
    q_h = Inches(1.3)
    gap = Inches(0.15)

    for num, label, question in questions:
        # Number block
        rect(s, Inches(0.7), q_y, Inches(0.7), q_h, fill=NAVY)
        txb(s, num, Inches(0.7), q_y, Inches(0.7), q_h,
            size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

        # Teal label
        rect(s, Inches(1.5), q_y, Inches(2.0), q_h, fill=TEAL)
        txb(s, label, Inches(1.5), q_y, Inches(2.0), q_h,
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

        # Question text
        txb(s, question, Inches(3.65), q_y + Inches(0.3), Inches(9.0), Inches(0.75),
            size=16, color=SLATE, font="Calibri")

        q_y += q_h + gap


def slide5_concepts(prs):
    """4 concept cards — teal left-accent strips."""
    s = blank_slide(prs)
    set_bg(s, WHITE)

    txb(s, "Core Concepts: Vertical & Horizontal Alignment",
        Inches(0.7), Inches(0.35), Inches(11.8), Inches(0.75),
        size=32, bold=True, color=NAVY, font="Calibri")
    txb(s, "Concept and Use",
        Inches(0.7), Inches(1.1), Inches(5), Inches(0.4),
        size=14, color=TEAL, bold=True, font="Calibri")

    concepts = [
        ("Continuity",   "Reiteration of major curriculum elements across educational levels."),
        ("Sequence",     "Progression of content from simple to complex — goes beyond continuity."),
        ("Integration",  "Horizontal relationship linking related subjects taught at the same grade."),
        ("Scope",        "The breadth and depth of content to be studied at a particular time."),
    ]

    card_y = Inches(1.65)
    card_h = Inches(1.08)
    gap    = Inches(0.2)

    for title, desc in concepts:
        # Card background
        rect(s, Inches(0.7), card_y, Inches(11.8), card_h,
             fill=LIGHT, line_color=RULE, line_width=Pt(0.5))
        # Teal left accent
        rect(s, Inches(0.7), card_y, Inches(0.1), card_h, fill=TEAL)
        # Title
        txb(s, title, Inches(1.0), card_y + Inches(0.15), Inches(2.5), Inches(0.4),
            size=17, bold=True, color=NAVY, font="Calibri")
        # Description
        txb(s, desc, Inches(1.0), card_y + Inches(0.58), Inches(10.5), Inches(0.42),
            size=13, color=SLATE, font="Calibri")

        card_y += card_h + gap


def concept_detail_slide(prs, title, points):
    """Generic concept detail slide — header + bullet points."""
    s = blank_slide(prs)
    set_bg(s, WHITE)

    # Teal left bar
    rect(s, Inches(0.7), Inches(0.4), Inches(0.1), Inches(6.6), fill=TEAL)

    txb(s, title,
        Inches(1.0), Inches(0.4), Inches(11.7), Inches(0.75),
        size=34, bold=True, color=NAVY, font="Calibri")

    # subtle rule
    hline(s, Inches(1.0), Inches(1.22), Inches(11.7), color=RULE, height=Inches(0.03))

    bullet_y = Inches(1.45)
    for pt in points:
        # bullet dot
        rect(s, Inches(1.0), bullet_y + Inches(0.14), Inches(0.1), Inches(0.1), fill=TEAL)
        # text
        txb(s, pt, Inches(1.25), bullet_y, Inches(11.4), Inches(0.85),
            size=15, color=SLATE, font="Calibri")
        bullet_y += Inches(1.0)


def slide6_continuity(prs):
    concept_detail_slide(prs, "Continuity", [
        "Reiteration of major elements of the curriculum across grade levels.",
        "Helps students stay motivated; gives priority to important knowledge and skills.",
        "Makes new content easier and more comprehensible by revisiting key ideas.",
        "There is no hard and fast rule for deciding the degree of continuity.",
        "Demands duplication of major curriculum components — not necessarily progress from simple to complex.",
    ])


def slide7_sequence(prs):
    concept_detail_slide(prs, "Sequence", [
        "Related to continuity, but goes beyond it — demands progress, not just repetition.",
        "Organising principles include: simple-to-complex, prerequisite learning, chronological order, whole-to-part, and increasing abstraction.",
        "Creates motivation within learners and builds specialised knowledge in a field.",
        "Demands not only duplication but progressive ordering from simple to complex.",
    ])


def slide8_example(prs):
    """Clean vertical flow chart for Continuity vs Sequence example."""
    s = blank_slide(prs)
    set_bg(s, WHITE)

    txb(s, "Continuity and Sequence: Example",
        Inches(0.7), Inches(0.35), Inches(11.8), Inches(0.7),
        size=32, bold=True, color=NAVY, font="Calibri")

    # ── LEFT SIDE: Continuity ────────────────────────────────────────
    label_x = Inches(0.7)
    txb(s, "CONTINUITY", label_x, Inches(1.3), Inches(2.2), Inches(0.4),
        size=13, bold=True, color=TEAL, font="Calibri")
    txb(s, "Same topic revisited\nacross grade levels",
        label_x, Inches(1.72), Inches(2.2), Inches(0.65),
        size=12, color=SLATE, font="Calibri")

    # Three "Addition" boxes stacked, connected
    cont_boxes = [
        (Inches(0.7), Inches(2.55), "Addition\n(Grade 1)"),
        (Inches(0.7), Inches(3.75), "Addition\n(Grade 2)"),
        (Inches(0.7), Inches(4.95), "Addition\n(Grade 3)"),
    ]
    box_w = Inches(2.3)
    box_h = Inches(0.9)
    for bx, by, bt in cont_boxes:
        rect(s, bx, by, box_w, box_h, fill=LIGHT, line_color=TEAL, line_width=Pt(1.5))
        txb(s, bt, bx, by, box_w, box_h,
            size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Calibri")
    # Connectors
    for i in range(len(cont_boxes) - 1):
        _, by, _ = cont_boxes[i]
        arrow_down(s, Inches(0.7) + box_w / 2, by + box_h + Inches(0.05),
                   Inches(0.8), color=MID)

    # ── RIGHT SIDE: Sequence ─────────────────────────────────────────
    seq_x = Inches(4.5)
    txb(s, "SEQUENCE", seq_x, Inches(1.3), Inches(8.0), Inches(0.4),
        size=13, bold=True, color=TEAL, font="Calibri")
    txb(s, "Content progresses from simple to complex",
        seq_x, Inches(1.72), Inches(8.0), Inches(0.45),
        size=12, color=SLATE, font="Calibri")

    seq_boxes = [
        ("Single-digit Addition",                "Recall of basic number bonds", "STEP 1"),
        ("Double-digit Addition",                "Place value understanding required", "STEP 2"),
        ("Double-digit Addition with Carry-over", "Regrouping and composing tens", "STEP 3"),
    ]
    colors = [LIGHT, LIGHT, LIGHT]
    accents = [RGBColor(0xCC, 0xE8, 0xE8), RGBColor(0xB2, 0xD8, 0xD8), TEAL_L]
    seq_ys = [Inches(2.55), Inches(3.9), Inches(5.2)]
    sbox_w = Inches(8.0)
    sbox_h = Inches(1.1)

    for (label, note, step), sy in zip(seq_boxes, seq_ys):
        rect(s, seq_x, sy, sbox_w, sbox_h,
             fill=LIGHT, line_color=TEAL, line_width=Pt(1.2))
        # Step badge
        rect(s, seq_x, sy, Inches(1.1), sbox_h, fill=TEAL)
        txb(s, step, seq_x, sy, Inches(1.1), sbox_h,
            size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")
        txb(s, label, seq_x + Inches(1.25), sy + Inches(0.1), sbox_w - Inches(1.4), Inches(0.45),
            size=16, bold=True, color=NAVY, font="Calibri")
        txb(s, note, seq_x + Inches(1.25), sy + Inches(0.58), sbox_w - Inches(1.4), Inches(0.42),
            size=12, color=SLATE, font="Calibri")

    # Upward arrows between seq boxes
    for i in range(len(seq_ys) - 1):
        mid_x = seq_x + sbox_w / 2
        arrow_down(s, mid_x, seq_ys[i] + sbox_h + Inches(0.05),
                   seq_ys[i + 1] - seq_ys[i] - sbox_h - Inches(0.05), color=TEAL)

    # Divider between left and right
    vline(s, Inches(3.5 - 0.02), Inches(1.3), Inches(5.8), color=RULE, width=Inches(0.03))


def slide9_scope(prs):
    concept_detail_slide(prs, "Scope", [
        "The breadth and depth of content to be studied in the curriculum at a particular time.",
        "Time constraints, common/core content, special learner needs, and integration all shape scope.",
        "Scope is related to the horizontal organisation of the curriculum.",
        "When curricularists decide what content to include and how intensively to teach it, they are defining scope.",
    ])


def slide10_scope_example(prs):
    """2-axis matrix: Breadth (x-axis) × Depth (y-axis)."""
    s = blank_slide(prs)
    set_bg(s, WHITE)

    txb(s, "Scope: An Example",
        Inches(0.7), Inches(0.35), Inches(11.8), Inches(0.7),
        size=32, bold=True, color=NAVY, font="Calibri")

    # ── Axis labels ──
    # X-axis label
    txb(s, "BREADTH  →  Topics covered at a given grade level",
        Inches(2.5), Inches(6.85), Inches(9.5), Inches(0.45),
        size=13, bold=True, color=TEAL, font="Calibri")

    # Y-axis label (rotated via a separate approach — we use a tall text box)
    txb(s, "D\nE\nP\nT\nH\n↑",
        Inches(0.35), Inches(2.1), Inches(0.4), Inches(4.2),
        size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER, font="Calibri")

    # Table setup
    subjects = ["Addition", "Subtraction", "Multiplication", "Division"]
    depths   = ["Single digit", "Double digit", "Three digit"]

    col_w  = Inches(2.2)
    row_h  = Inches(1.2)
    org_x  = Inches(1.05)
    org_y  = Inches(1.35)

    # Header row (subjects)
    for j, subj in enumerate(subjects):
        cx = org_x + j * col_w
        rect(s, cx, org_y, col_w - Inches(0.05), row_h * 0.6,
             fill=NAVY, line_color=None)
        txb(s, subj, cx, org_y, col_w - Inches(0.05), row_h * 0.6,
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

    # Depth rows
    for i, depth in enumerate(depths):
        ry = org_y + row_h * 0.6 + i * row_h

        # Row header
        rect(s, Inches(0.0), ry, org_x - Inches(0.05), row_h - Inches(0.05),
             fill=TEAL, line_color=None)
        txb(s, depth, Inches(0.0), ry, org_x - Inches(0.05), row_h - Inches(0.05),
            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri")

        for j in range(len(subjects)):
            cx = org_x + j * col_w
            # shade deeper rows darker to show depth
            fill_c = [LIGHT, RGBColor(0xE2, 0xEF, 0xEF), RGBColor(0xC8, 0xE6, 0xE6)][i]
            rect(s, cx, ry, col_w - Inches(0.05), row_h - Inches(0.05),
                 fill=fill_c, line_color=RULE, line_width=Pt(0.5))
            # Only subtraction column has explicit depth example
            cell_text = ""
            if j == 1:   # Subtraction
                cell_text = depth + "\nsubtraction"
            else:
                cell_text = depth
            txb(s, cell_text, cx + Inches(0.08), ry + Inches(0.15),
                col_w - Inches(0.2), row_h - Inches(0.25),
                size=12, color=SLATE, align=PP_ALIGN.CENTER, font="Calibri")

    # Note
    txb(s, "Breadth = number of topics at a grade level  |  Depth = complexity/intensity of a single topic",
        Inches(1.0), Inches(6.35), Inches(11.0), Inches(0.4),
        size=11, color=MID, align=PP_ALIGN.CENTER, font="Calibri", italic=True)


def slide11_integration(prs):
    concept_detail_slide(prs, "Integration", [
        "The horizontal relationship of curriculum elements (Tyler, 1949).",
        "Eliminates contradictions that arise between content in two or more related but separate subjects.",
        "Provides integrative learning — solving the problem of fragmentation and compartmentalisation.",
        "Refers to the relationship between different subjects taught at the same grade level.",
    ])


def slide12_closing(prs):
    """Minimalist closing slide."""
    s = blank_slide(prs)
    set_bg(s, NAVY)

    # Large teal circle accent (decorative)
    rect(s, Inches(10.2), Inches(-0.5), Inches(3.5), Inches(3.5), fill=TEAL)

    # Thank you
    txb(s, "Thank You",
        Inches(1.0), Inches(1.8), Inches(11.0), Inches(2.0),
        size=72, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri")

    # Rule
    rect(s, Inches(1.0), Inches(4.0), Inches(3.5), Inches(0.05), fill=TEAL)

    # Q&A
    txb(s, "Question and Answer Session",
        Inches(1.0), Inches(4.3), Inches(10.0), Inches(0.65),
        size=24, color=RGBColor(0xAD, 0xB5, 0xBD), font="Calibri")

    # Subtle footer
    txb(s, "Vertical and Horizontal Alignments in Curriculum  ·  Dhruba Prasad Niure, PhD",
        Inches(0.7), Inches(6.9), Inches(12.0), Inches(0.4),
        size=11, color=RGBColor(0x6C, 0x75, 0x7D), font="Calibri")


# ──────────────────────────────────────────────────────────────────────────────
#  MAIN
# ──────────────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()

    slide1_title(prs)
    slide2_curriculum(prs)
    slide3_alignment(prs)
    slide4_questions(prs)
    slide5_concepts(prs)
    slide6_continuity(prs)
    slide7_sequence(prs)
    slide8_example(prs)
    slide9_scope(prs)
    slide10_scope_example(prs)
    slide11_integration(prs)
    slide12_closing(prs)

    out = "/mnt/user-data/outputs/Vertical_Horizontal_Alignments_Redesigned.pptx"
    prs.save(out)
    print(f"Saved → {out}")

build()